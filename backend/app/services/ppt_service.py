import io, os, uuid, requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme definitions ─────────────────────────────────────────────────────────
THEMES = {
    "dark":    {"bg": RGBColor(0x0D, 0x10, 0x17), "primary": RGBColor(0x6C, 0x63, 0xFF),
                "text": RGBColor(0xFF, 0xFF, 0xFF), "accent": RGBColor(0x00, 0xD4, 0xFF),
                "sub": RGBColor(0xAA, 0xAA, 0xCC)},
    "light":   {"bg": RGBColor(0xFF, 0xFF, 0xFF), "primary": RGBColor(0x4F, 0x46, 0xE5),
                "text": RGBColor(0x1A, 0x1A, 0x2E), "accent": RGBColor(0x06, 0xB6, 0xD4),
                "sub": RGBColor(0x64, 0x64, 0x80)},
    "ocean":   {"bg": RGBColor(0x04, 0x2A, 0x4A), "primary": RGBColor(0x00, 0xB4, 0xD8),
                "text": RGBColor(0xF0, 0xF9, 0xFF), "accent": RGBColor(0x90, 0xE0, 0xEF),
                "sub": RGBColor(0x70, 0xB8, 0xD0)},
    "forest":  {"bg": RGBColor(0x08, 0x2A, 0x14), "primary": RGBColor(0x2D, 0x6A, 0x4F),
                "text": RGBColor(0xD8, 0xF3, 0xDC), "accent": RGBColor(0x52, 0xB7, 0x88),
                "sub": RGBColor(0x95, 0xD5, 0xB2)},
    "sunset":  {"bg": RGBColor(0x1A, 0x05, 0x33), "primary": RGBColor(0xFF, 0x63, 0x48),
                "text": RGBColor(0xFF, 0xF0, 0xF5), "accent": RGBColor(0xFF, 0xA5, 0x00),
                "sub": RGBColor(0xFF, 0x90, 0x80)},
    "minimal": {"bg": RGBColor(0xFA, 0xFA, 0xFA), "primary": RGBColor(0x22, 0x22, 0x22),
                "text": RGBColor(0x22, 0x22, 0x22), "accent": RGBColor(0x88, 0x88, 0x88),
                "sub": RGBColor(0x66, 0x66, 0x66)},
}

W = Inches(13.33)   # Widescreen 16:9
H = Inches(7.5)


def _blank_slide(prs: Presentation, theme: dict):
    layout = prs.slide_layouts[6]   # blank
    slide  = prs.slides.add_slide(layout)
    bg     = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = theme["bg"]
    return slide


def _add_rect(slide, l, t, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, text, l, t, w, h, color, size, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    return txb


def _add_image_from_url(slide, url: str, l, t, w, h):
    try:
        resp = requests.get(url, timeout=8)
        img_stream = io.BytesIO(resp.content)
        slide.shapes.add_picture(img_stream, l, t, w, h)
    except Exception as e:
        print(f"Image insert error: {e}")


def build_title_slide(slide, data: dict, theme: dict, image_url: str = None):
    # Full-bleed image with dark overlay
    if image_url:
        _add_image_from_url(slide, image_url, 0, 0, W, H)
        _add_rect(slide, 0, 0, W, H, theme["bg"])  # semi-overlay (solid for now)

    # Accent bar
    _add_rect(slide, 0, Inches(3.2), Inches(0.08), Inches(1.4), theme["primary"])

    # Title
    _add_text(slide, data.get("title", ""), Inches(0.5), Inches(2.5), Inches(12),
              Inches(1.4), theme["text"], 44, bold=True, align=PP_ALIGN.LEFT)

    # Subtitle
    if data.get("subtitle"):
        _add_text(slide, data["subtitle"], Inches(0.5), Inches(4.2), Inches(10),
                  Inches(0.9), theme["accent"], 22, align=PP_ALIGN.LEFT)


def build_content_slide(slide, data: dict, theme: dict, image_url: str = None):
    # Left panel accent
    _add_rect(slide, 0, 0, Inches(0.06), H, theme["primary"])

    # Title
    _add_text(slide, data.get("title", ""), Inches(0.4), Inches(0.3), Inches(12.5),
              Inches(1.0), theme["text"], 32, bold=True)

    # Divider
    _add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.03), theme["primary"])

    # Bullet points
    content_w = Inches(7.5) if image_url else Inches(12.5)
    points = data.get("points", [])
    for i, point in enumerate(points[:5]):
        y = Inches(1.5) + i * Inches(0.95)
        # Bullet dot
        _add_rect(slide, Inches(0.4), y + Inches(0.22), Inches(0.12), Inches(0.12), theme["accent"])
        _add_text(slide, point, Inches(0.7), y, content_w, Inches(0.85), theme["text"], 18)

    # Image on right
    if image_url:
        _add_image_from_url(slide, image_url, Inches(8.5), Inches(1.4), Inches(4.5), Inches(5.8))


def build_statistic_slide(slide, data: dict, theme: dict):
    # Big stat in center
    _add_rect(slide, 0, 0, W, H, theme["bg"])
    _add_text(slide, data.get("title", ""), Inches(0.5), Inches(0.3), W, Inches(0.9),
              theme["sub"], 24, bold=False, align=PP_ALIGN.CENTER)

    stat = data.get("stat", "")
    _add_text(slide, stat, Inches(0.5), Inches(1.2), W, Inches(2.5),
              theme["primary"], 96, bold=True, align=PP_ALIGN.CENTER)

    label = data.get("stat_label", "")
    _add_text(slide, label, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
              theme["accent"], 22, align=PP_ALIGN.CENTER)

    _add_rect(slide, Inches(4), Inches(4.6), Inches(5.33), Inches(0.04), theme["primary"])

    points = data.get("points", [])
    for i, pt in enumerate(points[:3]):
        _add_text(slide, f"• {pt}", Inches(0.5), Inches(4.9) + i * Inches(0.7),
                  W, Inches(0.6), theme["sub"], 16, align=PP_ALIGN.CENTER)


def build_agenda_slide(slide, data: dict, theme: dict):
    _add_rect(slide, 0, 0, Inches(4), H, theme["primary"])
    _add_text(slide, "AGENDA", Inches(0.3), Inches(0.4), Inches(3.5), Inches(0.8),
              theme["text"], 28, bold=True)

    points = data.get("points", [])
    for i, point in enumerate(points[:6]):
        y = Inches(1.2) + i * Inches(0.9)
        _add_rect(slide, Inches(0.3), y, Inches(0.4), Inches(0.4),
                  RGBColor(0xFF, 0xFF, 0xFF) if True else theme["text"])
        _add_text(slide, str(i + 1), Inches(0.35), y + Inches(0.02),
                  Inches(0.3), Inches(0.4), theme["bg"], 16, bold=True, align=PP_ALIGN.CENTER)
        _add_text(slide, point, Inches(0.9), y, Inches(2.8), Inches(0.45), theme["text"], 16)

    _add_text(slide, data.get("title", "Agenda"), Inches(4.5), Inches(0.4),
              Inches(8.5), Inches(1.2), theme["text"], 38, bold=True)
    _add_rect(slide, Inches(4.5), Inches(1.5), Inches(8), Inches(0.04), theme["accent"])


def build_conclusion_slide(slide, data: dict, theme: dict, image_url: str = None):
    if image_url:
        _add_image_from_url(slide, image_url, 0, 0, W, H)
    _add_rect(slide, 0, 0, W, H, RGBColor(0x08, 0x08, 0x12))

    _add_rect(slide, Inches(0.5), Inches(0.8), Inches(1.5), Inches(0.06), theme["primary"])
    _add_text(slide, "KEY TAKEAWAYS", Inches(0.5), Inches(1.0), Inches(12),
              Inches(0.7), theme["accent"], 18, bold=True)
    _add_text(slide, data.get("title", "Conclusion"), Inches(0.5), Inches(1.6),
              Inches(12), Inches(1.0), theme["text"], 36, bold=True)

    points = data.get("points", [])
    for i, pt in enumerate(points[:4]):
        y = Inches(2.8) + i * Inches(0.9)
        _add_rect(slide, Inches(0.5), y + Inches(0.2), Inches(0.3), Inches(0.3), theme["primary"])
        _add_text(slide, pt, Inches(1.1), y, Inches(11.5), Inches(0.8), theme["text"], 18)


def generate_pptx(presentation_data: dict, theme_name: str, image_map: dict) -> str:
    theme = THEMES.get(theme_name, THEMES["dark"])
    prs   = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for i, slide_data in enumerate(presentation_data["slides"]):
        slide_type = slide_data.get("type", "content")
        img_url    = image_map.get(i)
        slide      = _blank_slide(prs, theme)

        if slide_type == "title":
            build_title_slide(slide, slide_data, theme, img_url)
        elif slide_type == "agenda":
            build_agenda_slide(slide, slide_data, theme)
        elif slide_type == "statistic":
            build_statistic_slide(slide, slide_data, theme)
        elif slide_type == "conclusion":
            build_conclusion_slide(slide, slide_data, theme, img_url)
        else:   # content / comparison
            build_content_slide(slide, slide_data, theme, img_url)

        # Slide number (footer)
        _add_text(slide, str(i + 1), Inches(12.8), Inches(7.1), Inches(0.4),
                  Inches(0.3), theme["sub"], 10, align=PP_ALIGN.RIGHT)

    file_id  = str(uuid.uuid4())[:8]
    filename = f"app/generated/slideai_{file_id}.pptx"
    prs.save(filename)
    return filename, file_id
